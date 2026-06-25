"""Generate the ForgeAI presentation deck (PPTX).

Run:  uv run --with python-pptx python scripts/make_deck.py
Output: ForgeAI-Presentation.pptx in the repo root.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

# ---- palette ---------------------------------------------------------------
BG = RGBColor(0x0B, 0x0E, 0x14)      # near-black
PANEL = RGBColor(0x15, 0x1A, 0x24)   # dark panel
INK = RGBColor(0xF2, 0xF4, 0xF8)     # near-white
MUTE = RGBColor(0x9A, 0xA4, 0xB2)    # muted grey
ACCENT = RGBColor(0x4F, 0x9C, 0xF5)  # blue
GREEN = RGBColor(0x3F, 0xC9, 0x7A)
AMBER = RGBColor(0xF2, 0xC1, 0x4E)

# 16:9
prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h, fill=None, line=None):
    shp = s.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, space=6, font="Calibri"):
    tb = s.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, size, color, bold)]
    for i, item in enumerate(runs):
        t, sz, col, bd = (item if len(item) == 4 else (*item, False))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = t
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.bold = bd
        r.font.name = font
    return tb


def bullets(s, x, y, w, h, items, size=16, gap=10):
    tb = s.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        lead, rest = item if isinstance(item, tuple) else (item, "")
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "▸ "
        r.font.size = Pt(size)
        r.font.color.rgb = ACCENT
        r.font.bold = True
        r2 = p.add_run()
        r2.text = lead
        r2.font.size = Pt(size)
        r2.font.color.rgb = INK
        r2.font.bold = True
        r2.font.name = "Calibri"
        if rest:
            r3 = p.add_run()
            r3.text = " — " + rest
            r3.font.size = Pt(size)
            r3.font.color.rgb = MUTE
            r3.font.name = "Calibri"
    return tb


def kicker(s, label):
    text(s, 700000, 520000, 8000000, 400000, label.upper(), size=13, color=ACCENT, bold=True)


def title(s, t, y=820000):
    text(s, 700000, y, 10800000, 900000, t, size=34, color=INK, bold=True)


M = 700000  # left margin

# ===========================================================================
# 1 — TITLE
# ===========================================================================
s = slide()
box(s, 0, 2550000, SW, Emu(14000), fill=ACCENT)
text(s, M, 1700000, 10800000, 900000, "ForgeAI", size=66, color=INK, bold=True)
text(s, M, 2750000, 10800000, 600000,
     "An autonomous AI engineering team — that measures and improves itself.",
     size=22, color=MUTE)
text(s, M, 3600000, 11000000, 1400000, [
    ("A multi-agent platform that plans, codes, tests, reviews, and ships —", 17, INK, False),
    ("then scores every run and gets measurably better over time.", 17, INK, False),
], size=17)
text(s, M, 5750000, 11000000, 500000,
     "13 phases · 330 tests · fully offline-testable · runs on local models",
     size=14, color=MUTE)

# ===========================================================================
# 2 — THE PROBLEM
# ===========================================================================
s = slide()
kicker(s, "The problem")
title(s, "Most AI coding tools are one model, one prompt")
# left card
box(s, M, 2050000, 5100000, 3400000, fill=PANEL, line=RGBColor(0x2A,0x31,0x40))
text(s, M+300000, 2300000, 4600000, 500000, "TODAY", size=14, color=MUTE, bold=True)
text(s, M+300000, 2900000, 4600000, 1500000, [
    ("User", 22, INK, True), ("↓", 20, MUTE, False),
    ("LLM", 22, INK, True), ("↓", 20, MUTE, False),
    ("Answer", 22, INK, True),
], align=PP_ALIGN.CENTER)
text(s, M+300000, 4900000, 4600000, 500000,
     "One generalist. No plan, no tests, no review, no memory of what worked.",
     size=14, color=MUTE)
# right card
box(s, 6600000, 2050000, 5000000, 3400000, fill=PANEL, line=ACCENT)
text(s, 6900000, 2300000, 4400000, 500000, "FORGEAI", size=14, color=ACCENT, bold=True)
text(s, 6900000, 2850000, 4500000, 2200000, [
    ("A coordinated team of specialists", 18, INK, True),
    ("Manager → Planner → Researcher → Coder", 15, MUTE, False),
    ("→ Tester → Reviewer → DevOps → PR", 15, MUTE, False),
    ("", 8, MUTE, False),
    ("…then every run is scored, stored, and", 15, GREEN, False),
    ("compared — so it improves without code changes.", 15, GREEN, False),
], size=15)

# ===========================================================================
# 3 — WHAT IT IS
# ===========================================================================
s = slide()
kicker(s, "What it is")
title(s, "A team of AI engineers — and a self-improving system")
bullets(s, M, 1950000, 10900000, 3600000, [
    ("Multi-agent autonomous engineering", "10 specialist agents, each with one job, coordinated by an explicit workflow graph over one shared state."),
    ("Natural language → reviewed pull request", "describe a task; watch the team plan, code, test, and review it live inside a sandbox."),
    ("A self-improving layer on top", "every run is measured, stored, and compared — prompt versions, benchmarks, failure memory, learning loops."),
    ("Runs and is fully tested offline", "every external dependency (LLM, GitHub, datastores) has a deterministic fake. 330 tests pass in seconds."),
], size=18, gap=16)

# ===========================================================================
# 4 — WHY I BUILT IT THIS WAY
# ===========================================================================
s = slide()
kicker(s, "Why — the design decisions")
title(s, "Five choices that shape the whole system")
cards = [
    ("A team, not a model", "Real software is built by specialists. Splitting work makes each step inspectable, testable, and debuggable."),
    ("The graph sequences agents", "Agents never call each other. Order, branching, and retries are explicit and visible (loose coupling)."),
    ("Provider abstraction + fakes", "Same code offline (echo, fake GitHub) and in prod (Ollama, real GitHub). Going live is just config."),
    ("Human-gated writes", "PRs, prompt promotions, workflow changes are proposed then approved. Never silently mutates production."),
    ("Measure before optimizing", "Build the evaluation substrate first; learning loops stay advisory until there's real data to act on."),
]
x, y, w, h, gap = M, 1950000, 3450000, 1750000, 280000
for i, (t, d) in enumerate(cards):
    col = i % 3
    row = i // 3
    cx = x + col * (w + gap)
    cy = y + row * (h + gap)
    box(s, cx, cy, w, h, fill=PANEL, line=RGBColor(0x2A,0x31,0x40))
    text(s, cx+200000, cy+170000, w-400000, h-300000, [
        (f"{i+1}. {t}", 15, ACCENT, True),
        (d, 12.5, MUTE, False),
    ], size=13, space=5)

# ===========================================================================
# 5 — HOW IT WORKS (PIPELINE)
# ===========================================================================
s = slide()
kicker(s, "How it works")
title(s, "One request flows through an explicit workflow graph")
flow = ["Manager\n(intake)", "Planner", "Researcher", "Memory", "Coder",
        "Execute\n(sandbox)", "Tests", "Review"]
bx, by, bw, bh, g = M, 2150000, 1230000, 900000, 130000
for i, label in enumerate(flow):
    cx = bx + i * (bw + g)
    box(s, cx, by, bw, bh, fill=PANEL, line=ACCENT)
    text(s, cx, by, bw, bh, label, size=12.5, color=INK, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, M, 3300000, 11000000, 600000,
     "Review →  approved → Git (proposes a gated PR) → final.   Not approved → Reflection → Coder (retry loop).",
     size=14, color=MUTE)
box(s, M, 4050000, 10900000, 1850000, fill=PANEL, line=RGBColor(0x2A,0x31,0x40))
text(s, M+300000, 4230000, 10400000, 1500000, [
    ("Around the graph:", 15, INK, True),
    ("Observability — every node emits events to a live timeline, metrics, and WebSocket.", 13.5, MUTE, False),
    ("Sandbox — generated code runs in isolated Docker; destructive ops withheld.", 13.5, MUTE, False),
    ("Memory + RAG — past tasks and context retrieved from Qdrant for new requests.", 13.5, MUTE, False),
    ("GitHub — commits on a local clone; PRs opened only after human approval.", 13.5, MUTE, False),
], size=13.5, space=6)

# ===========================================================================
# 6 — SELF-IMPROVEMENT (PHASE 12)
# ===========================================================================
s = slide()
kicker(s, "The differentiator — Phase 12")
title(s, "It measures itself, then improves — behind approval gates")
left = [
    ("Evaluation Engine", "scores every run via a versioned rubric"),
    ("Performance Database", "durable, comparable run stats (derived, can't drift)"),
    ("Prompt Versioning", "append-only; every run records the version it used"),
    ("Failure Knowledge Base", "Error → Fix → Store → Reuse; self-correcting"),
    ("Benchmark Suite", "versioned scenarios scored per release"),
]
right = [
    ("Multi-Agent Debate", "N planners compete; a judge picks the winner"),
    ("Dynamic Selection", "routes by task type behind a swappable strategy"),
    ("Agent Analytics", "dashboard: deltas, prompt comparison, trend"),
    ("Learning loops", "A/B promote, workflow-opt, PR-outcome — gated"),
    ("Agent Marketplace", "register + discover contract-checked plugins"),
]
bullets(s, M, 1950000, 5450000, 3600000, left, size=14.5, gap=11)
bullets(s, 6450000, 1950000, 5300000, 3600000, right, size=14.5, gap=11)
box(s, M, 5650000, 10900000, 700000, fill=PANEL, line=GREEN)
text(s, M, 5650000, 10900000, 700000,
     "Proven live:  type a task  →  team executes  →  run scored  →  persisted to PostgreSQL  →  shown on the dashboard.",
     size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# 6.5 — PROJECTS & FIRST-RUN (PHASE 13)
# ===========================================================================
s = slide()
kicker(s, "From engine to product — Phase 13")
title(s, "A front door: projects, bootstrap, first-minute wow")
p13 = [
    ("First-class Projects", "each owns a real workspace directory; CRUD within the existing tenancy + RBAC"),
    ("Runs bind to a project", "/agents/run targets its path and writes generated files there — a real project on disk"),
    ("Bootstrap from nothing", "versioned starters (empty, or FastAPI + JWT + Postgres + Docker + tests) scaffold instantly, offline"),
    ("Onboarding flow", "sign in → project chooser (Create New / Open Existing) → workspace bound to the project"),
    ("The first-minute wow", "fresh account → pick a starter → watch the team build it live; works offline (MODEL_PROVIDER=echo)"),
]
bullets(s, M, 1950000, 10900000, 3600000, p13, size=15.5, gap=13)
box(s, M, 5650000, 10900000, 700000, fill=PANEL, line=ACCENT)
text(s, M, 5650000, 10900000, 700000,
     'Closed the critique head-on: "to which project?", "no bootstrap", "no sub-30s wow" — all resolved.',
     size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# 7 — BENEFITS
# ===========================================================================
s = slide()
kicker(s, "Benefits")
title(s, "Why this approach pays off")
cards = [
    ("Inspectable & debuggable", "Each agent does one job; the timeline shows exactly who did what, when."),
    ("Measurably better over time", "Versioned rubric + benchmarks turn 'did it help?' into a number, not a vibe."),
    ("Safe by construction", "Sandboxed execution, path validation, JWT auth, human-gated external writes."),
    ("Runs fully local", "Ollama models, no API keys, no data leaving your machine. Swap by config."),
    ("Reproducible", "Deterministic fakes → the whole system tests offline; great for CI & demos."),
    ("Extensible", "Provider abstractions + plugin marketplace; add agents/models without rewrites."),
]
x, y, w, h, gap = M, 1950000, 3450000, 1700000, 280000
for i, (t, d) in enumerate(cards):
    col, row = i % 3, i // 3
    cx = x + col * (w + gap)
    cy = y + row * (h + gap)
    box(s, cx, cy, w, h, fill=PANEL, line=RGBColor(0x2A,0x31,0x40))
    text(s, cx+200000, cy+200000, w-400000, h-350000, [
        (t, 15, GREEN, True), (d, 13, MUTE, False),
    ], size=13, space=6)

# ===========================================================================
# 8 — APPLICATIONS
# ===========================================================================
s = slide()
kicker(s, "Applications")
title(s, "Two ways to use it")
box(s, M, 1950000, 5300000, 3700000, fill=PANEL, line=ACCENT)
text(s, M+300000, 2150000, 4800000, 500000, "AS AN ENGINEERING TEAM", size=14, color=ACCENT, bold=True)
bullets(s, M+300000, 2700000, 4800000, 2800000, [
    ("Developers", "offload features, fixes, refactors, tests"),
    ("Startups", "ship CRUD APIs & internal tooling fast"),
    ("Companies", "internal dashboards on local models, RBAC, gated PRs"),
    ("Students", "a transparent reference for how it all fits together"),
], size=14, gap=12)
box(s, 6450000, 1950000, 5300000, 3700000, fill=PANEL, line=GREEN)
text(s, 6750000, 2150000, 4800000, 500000, "AS A RESEARCH SUBSTRATE", size=14, color=GREEN, bold=True)
bullets(s, 6750000, 2700000, 4800000, 2800000, [
    ("Agent evaluation", "rubric + benchmarks + performance DB"),
    ("Prompt engineering at scale", "versioned prompts as A/B experiments"),
    ("Self-improving agents", "reflection, debate, learning-from-outcomes"),
    ("Plugin ecosystem", "permissioned, contract-checked agent marketplace"),
], size=14, gap=12)

# ===========================================================================
# 9 — STACK & STATUS / CLOSING
# ===========================================================================
s = slide()
kicker(s, "Stack & status")
title(s, "Built, tested, and running")
bullets(s, M, 1950000, 11000000, 2600000, [
    ("Frontend", "Next.js · React · TypeScript · Tailwind"),
    ("Backend", "FastAPI · LangGraph · SQLAlchemy (async) · WebSockets"),
    ("Data", "PostgreSQL · Redis · Qdrant"),
    ("AI", "Ollama (local) · model-router abstraction · LangGraph"),
    ("Infra", "Docker · Docker Compose"),
], size=16, gap=11)
box(s, M, 4750000, 10900000, 1150000, fill=PANEL, line=GREEN)
text(s, M, 4850000, 10900000, 1000000, [
    ("All 13 phases complete  ·  330 tests passing  ·  26 ADRs  ·  fully documented", 16, GREEN, True),
    ("From a static pipeline to a self-improving, approval-gated engineering team.", 14, MUTE, False),
], size=15, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=6)

prs.save("ForgeAI-Presentation.pptx")
print("Saved ForgeAI-Presentation.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
